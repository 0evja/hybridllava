from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from copy import deepcopy
from lxml import etree
import os

pptx_path = '/home/hechen/zms/MLLM_Factory/HiDe-LLaVA/中期答辩汇报PPT.pptx'
out_path = '/home/hechen/zms/MLLM_Factory/HiDe-LLaVA/中期答辩汇报PPT_improved.pptx'
charm_dir = '/home/hechen/zms/MLLM_Factory/HiDe-LLaVA/charm'

prs = Presentation(pptx_path)
SW = prs.slide_width
SH = prs.slide_height

# 颜色
BLUE = RGBColor(0x44, 0x72, 0xC4)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
L_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
L_ORANGE = RGBColor(0xFD, 0xE8, 0xD0)
L_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
L_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
L_RED = RGBColor(0xFD, 0xED, 0xEC)
L_GRAY = RGBColor(0xF2, 0xF3, 0xF4)

def clear_slide(slide):
    for sp in list(slide.shapes):
        slide.shapes._spTree.remove(sp._element)

def title_bar(slide, text):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Emu(731520))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()
    tx = slide.shapes.add_textbox(Emu(457200), Emu(137160), Emu(8229600), Emu(548640))
    r = tx.text_frame.paragraphs[0].add_run()
    r.text = text; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE

def rbox(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(1.5)
    else: s.line.fill.background()
    return s

def txt(slide, l, t, w, h, lines, sz=18, clr=DARK, bold=False, sp=Pt(4), align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tx.text_frame; tf.word_wrap = True
    if isinstance(lines, str): lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line; r.font.size = Pt(sz)
        r.font.color.rgb = clr; r.font.bold = bold; p.space_after = sp; p.alignment = align
    return tx

def img(slide, path, l, t, w=None, h=None):
    kw = {}
    if w: kw['width'] = Emu(w)
    if h: kw['height'] = Emu(h)
    slide.shapes.add_picture(path, Emu(l), Emu(t), **kw)

def add_blank_slide(prs, insert_index):
    """在指定位置插入空白slide"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    # 移动到指定位置
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_slide = slides[-1]
    xml_slides.remove(new_slide)
    xml_slides.insert(insert_index, new_slide)
    return slide

# ============================================================
# Slide 9: 分隔页 — 改文字
# ============================================================
slide9 = prs.slides[8]
for shape in slide9.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if '研究内容' in p.text:
                for run in p.runs:
                    run.text = run.text.replace('研究内容与解决方案', '本研究的方案')
                if not p.runs:
                    p.text = '本研究的方案'

# ============================================================
# 新增 Slide: 本研究的目标（插入到原slide10之前，即index=9）
# ============================================================
new_slide = add_blank_slide(prs, 9)
title_bar(new_slide, '本研究的目标')

# 题目
txt(new_slide, 500000, 900000, 8100000, 450000,
    '课题名称：基于参数隔离的多模态大模型持续学习研究',
    sz=20, bold=True, clr=BLUE)

# 总目标
txt(new_slide, 500000, 1500000, 8100000, 400000,
    '研究目标', sz=20, bold=True, clr=DARK)

txt(new_slide, 700000, 1950000, 7700000, 900000, [
    '针对多模态大模型在连续指令微调中的灾难性遗忘问题，',
    '研究基于参数隔离的持续学习方法，构建一种在不显著增加',
    '计算开销的前提下，有效缓解新旧任务知识冲突的学习框架。',
], sz=17, clr=DARK, sp=Pt(6))

# 三个具体目标
txt(new_slide, 500000, 3000000, 8100000, 400000,
    '具体研究内容', sz=20, bold=True, clr=DARK)

goals = [
    ('1', '分层参数隔离架构研究',
     '探索如何将模型底层通用表征与顶层任务特异知识进行物理隔离，通过分层解耦的LoRA专家结构实现任务间参数的有效隔离。'),
    ('2', '任务自适应融合机制研究',
     '研究如何利用任务嵌入动态生成各层LoRA融合权重，替代固定等权策略，使参数隔离与融合具备任务自适应能力。'),
    ('3', '任务提示增强与记忆回顾策略研究',
     '探索通过可学习任务提示提供显式任务标识，并结合少量样本回放机制，从多个维度协同提升模型的抗遗忘能力。'),
]

for i, (num, title, desc) in enumerate(goals):
    y = 3500000 + i * 1050000
    # 编号圆
    s = new_slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(550000), Emu(y), Emu(400000), Emu(400000))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()
    tf = s.text_frame; r = tf.paragraphs[0].add_run()
    r.text = num; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    txt(new_slide, 1100000, y - 20000, 7200000, 350000,
        title, sz=18, bold=True, clr=BLUE)
    txt(new_slide, 1100000, y + 350000, 7200000, 600000,
        desc, sz=15, clr=GRAY, sp=Pt(3))

# ============================================================
# 原Slide 10 → 现在是 index 10: 现有方法的不足
# ============================================================
slide10 = prs.slides[10]
clear_slide(slide10)
title_bar(slide10, '现有方法的优化空间')

txt(slide10, 500000, 900000, 8100000, 500000,
    '通过对现有方法的梳理，我们发现以HiDe-LLaVA为代表的分层解耦方法虽然取得了显著进展，但仍存在以下优化空间：',
    sz=17, clr=DARK)

# 三个优化空间，用左侧竖条+白底卡片，更清爽
items = [
    ('优化空间一', '底层LoRA融合策略过于简单',
     '底层（Layer 0-30）对所有已学任务的LoRA采用固定等权融合（ε=1.0），无法区分不同任务对各层的贡献差异，缺乏自适应能力。',
     BLUE),
    ('优化空间二', '缺乏显式的任务身份标识',
     '模型仅通过CLIP特征与锚点的余弦相似度进行任务路由，缺少对任务身份的显式编码，路由精度受限于CLIP特征空间的区分能力。',
     ORANGE),
    ('优化空间三', '无记忆回顾机制',
     '训练新任务时完全不回顾旧任务数据，仅依赖架构隔离抵抗遗忘，对长序列任务的知识保持仍有不足。',
     PURPLE),
]

for i, (label, title, desc, color) in enumerate(items):
    y = 1600000 + i * 1600000
    # 左侧竖条
    bar = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Emu(500000), Emu(y), Emu(80000), Emu(1300000))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    # 标题
    txt(slide10, 750000, y + 50000, 7600000, 350000,
        label, sz=14, clr=color, bold=True)
    txt(slide10, 750000, y + 300000, 7600000, 400000,
        title, sz=19, clr=DARK, bold=True)
    # 描述
    txt(slide10, 750000, y + 720000, 7600000, 550000,
        desc, sz=15, clr=GRAY)

# ============================================================
# 原Slide 11 → 现在是 index 11: 解决方案
# ============================================================
slide11 = prs.slides[11]
clear_slide(slide11)
title_bar(slide11, '本研究的解决方案')

txt(slide11, 500000, 850000, 8100000, 400000,
    '针对三个优化空间，提出对应的改进策略：',
    sz=18, bold=True, clr=DARK)

# 三个方案，同样用左侧竖条风格
solutions = [
    ('针对优化空间一', '自适应融合系数',
     '用可学习的任务嵌入生成底层每层的LoRA融合权重，替代固定等权策略。',
     '固定: ΔW = Σ 1.0·BᵢAᵢ  →  自适应: ΔW = Σ wᵢ·BᵢAᵢ',
     BLUE),
    ('针对优化空间二', '可学习的任务提示',
     '为每个任务引入Task Prompt [10×4096]，注入输入序列，提供显式任务身份标识。',
     '输入序列: [sys] [image] [prompt] [text]',
     ORANGE),
    ('针对优化空间三', '经验回放（计划中）',
     '为每个已学任务保留5-10%代表性样本，训练新任务时混合回放，提供旧任务梯度信号。',
     '→ 后续v2.0实现',
     PURPLE),
]

for i, (label, title, desc, note, color) in enumerate(solutions):
    y = 1400000 + i * 1650000
    # 左侧竖条
    bar = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Emu(500000), Emu(y), Emu(80000), Emu(1350000))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    # 标题
    txt(slide11, 750000, y + 50000, 7600000, 300000,
        label, sz=14, clr=color, bold=True)
    txt(slide11, 750000, y + 300000, 7600000, 400000,
        title, sz=19, clr=DARK, bold=True)
    # 描述
    txt(slide11, 750000, y + 720000, 7600000, 350000,
        desc, sz=15, clr=GRAY)
    # 补充说明
    txt(slide11, 750000, y + 1050000, 7600000, 300000,
        note, sz=14, clr=color)

# ============================================================
# 原Slide 12 → 现在是 index 12: 分隔页 — 改编号
# ============================================================
slide12 = prs.slides[12]
for shape in slide12.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.text.strip() == '03':
                for run in p.runs:
                    run.text = '04'
                if not p.runs:
                    p.text = '04'

# ============================================================
# 原Slide 12 → 现在是 index 13: 综述调研与阅读（保留图片，重组文字）
# ============================================================
slide13 = prs.slides[13]
# 只删除非图片元素，保留3张调研图片
for sp in list(slide13.shapes):
    if sp.shape_type != MSO_SHAPE_TYPE.PICTURE:
        slide13.shapes._spTree.remove(sp._element)

# 重建标题栏
title_bar(slide13, '综述调研与阅读')

# 左栏：综述论文
txt(slide13, 500000, 900000, 4000000, 350000,
    '综述论文', sz=19, bold=True, clr=BLUE)

surveys = [
    ('1. Continual Learning for VLMs: A Survey and',
     '   Taxonomy Beyond Forgetting (Liu et al., 2025)',
     '→ VLM持续学习方法的深度分类总结'),
    ('2. When CL Meets MLLM: A Survey',
     '   (Huo & Tang, 2025)',
     '→ 持续学习与多模态融合方向的全面梳理'),
    ('3. Parameter-efficient Fine-tuning in LLMs:',
     '   A Survey of Methodologies',
     '→ 参数高效微调技术的系统性综述'),
]

for i, (t1, t2, arrow) in enumerate(surveys):
    y = 1350000 + i * 1050000
    # 左侧小竖条
    bar = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Emu(500000), Emu(y), Emu(60000), Emu(850000))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    txt(slide13, 700000, y + 30000, 4200000, 550000,
        [t1, t2], sz=13, clr=DARK, sp=Pt(2))
    txt(slide13, 700000, y + 550000, 4200000, 300000,
        arrow, sz=13, clr=BLUE, bold=True, sp=Pt(2))

# 右栏：核心方法论文
txt(slide13, 5200000, 900000, 3600000, 350000,
    '核心方法论文', sz=19, bold=True, clr=ORANGE)

methods = [
    ('经典持续学习', 'EWC、iCaRL、LwF'),
    ('参数高效微调', 'LoRA、MoE-LoRA'),
    ('多模态持续学习', 'HiDe-LLaVA (ACL 2025)'),
]

for i, (cat, papers) in enumerate(methods):
    y = 1350000 + i * 700000
    rbox(slide13, 5200000, y, 3500000, 580000, L_ORANGE, ORANGE)
    txt(slide13, 5380000, y + 80000, 3200000, 250000,
        cat, sz=15, bold=True, clr=ORANGE)
    txt(slide13, 5380000, y + 340000, 3200000, 250000,
        papers, sz=14, clr=DARK)

# 底部说明（图片在 y≈5100000，在其上方加文字）
txt(slide13, 500000, 4650000, 8100000, 350000,
    '调研过程中的文献梳理与分析', sz=15, clr=GRAY)

# ============================================================
# 原Slide 13 → 现在是 index 14: 基座模型部署与微调验证
# ============================================================
slide14 = prs.slides[14]
# 只删除非图片元素，保留右上角LLaVA架构图
for sp in list(slide14.shapes):
    if sp.shape_type != MSO_SHAPE_TYPE.PICTURE:
        slide14.shapes._spTree.remove(sp._element)

# 把LLaVA架构图移到右侧中部，放大使其清晰可见
for sp in slide14.shapes:
    if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
        sp.left = Emu(5100000)
        sp.top = Emu(850000)
        sp.width = Emu(3800000)
        sp.height = Emu(1330000)

title_bar(slide14, 'LLaVA 基座模型部署与微调验证')

# 左侧：说明文字（不覆盖右侧图片区域）
txt(slide14, 500000, 900000, 4400000, 500000,
    '在正式开展持续学习实验前，先对基座模型\n进行了完整的部署与微调验证。',
    sz=16, clr=GRAY)

# 基座模型信息（图片下方，全宽）
txt(slide14, 500000, 2350000, 8100000, 350000,
    '基座模型配置', sz=19, bold=True, clr=BLUE)

rbox(slide14, 500000, 2750000, 8100000, 750000, L_BLUE, BLUE)
txt(slide14, 680000, 2830000, 7700000, 600000, [
    'LLaVA-v1.5-7b：语言解码器 Vicuna-7B-v1.5 + 视觉编码器 CLIP-ViT-Large-patch14-336',
    '训练数据：COCO (Common Objects in Context) + Visual Genome (VG)',
], sz=15, clr=DARK, sp=Pt(8))

# 两种微调方式
txt(slide14, 500000, 3700000, 8100000, 350000,
    '两种微调方式对比验证', sz=19, bold=True, clr=BLUE)

# 方式1
rbox(slide14, 500000, 4100000, 3800000, 1400000, L_GREEN, GREEN)
txt(slide14, 680000, 4180000, 3400000, 350000,
    '方式一：仅微调 mm-projector', sz=16, bold=True, clr=GREEN)
txt(slide14, 680000, 4550000, 3400000, 850000, [
    '只更新视觉-语言投影层参数',
    '冻结 LLM 和视觉编码器',
    '→ 低成本适配，但表征调整有限',
], sz=14, clr=DARK, sp=Pt(6))

# 方式2
rbox(slide14, 4700000, 4100000, 3900000, 1400000, L_ORANGE, ORANGE)
txt(slide14, 4880000, 4180000, 3500000, 350000,
    '方式二：基于 LoRA 微调', sz=16, bold=True, clr=ORANGE)
txt(slide14, 4880000, 4550000, 3500000, 850000, [
    '在 LLM 层注入低秩适配矩阵',
    '参数高效，训练显存友好',
    '→ 更深层次的知识适配',
], sz=14, clr=DARK, sp=Pt(6))

# 结论
txt(slide14, 500000, 5700000, 8100000, 800000, [
    '→ 确认了两种微调粒度下模型行为的差异',
    '→ 验证了工程流程的正确性，为后续模块化隔离实验提供了可靠基础',
], sz=16, bold=True, clr=GREEN, sp=Pt(6))

# ============================================================
# 原Slide 15 → 现在是 index 15: 基线选择与复现动机
# ============================================================
slide15 = prs.slides[15]
clear_slide(slide15)
title_bar(slide15, '基线选择与复现动机')

txt(slide15, 500000, 900000, 8100000, 400000,
    '为什么选择 HiDe-LLaVA 作为基线？', sz=20, bold=True, clr=BLUE)

rbox(slide15, 400000, 1500000, 3900000, 2200000, L_BLUE, BLUE)
txt(slide15, 550000, 1600000, 3600000, 400000,
    'HiDe-LLaVA 的优势', sz=18, bold=True, clr=BLUE)
txt(slide15, 550000, 2050000, 3600000, 1600000, [
    '✓ ACL 2025，当前SOTA方法',
    '✓ 分层解耦思想新颖',
    '✓ 无需额外存储开销',
    '✓ 代码开源，可复现',
], sz=16, clr=DARK, sp=Pt(8))

rbox(slide15, 4700000, 1500000, 3900000, 2200000, L_RED, RED)
txt(slide15, 4850000, 1600000, 3600000, 400000,
    '我们发现的改进空间', sz=18, bold=True, clr=RED)
txt(slide15, 4850000, 2050000, 3600000, 1600000, [
    '✗ 底层融合固定等权',
    '✗ 缺乏任务身份标识',
    '✗ 无记忆回顾机制',
    '✗ 路由精度受限',
], sz=16, clr=DARK, sp=Pt(8))

txt(slide15, 500000, 4000000, 8100000, 400000,
    '→ 复现HiDe-LLaVA，验证基线性能，为后续改进提供对照基准',
    sz=18, bold=True, clr=GREEN)

rbox(slide15, 400000, 4600000, 8300000, 1800000, L_GRAY, GRAY)
txt(slide15, 550000, 4700000, 7900000, 400000,
    '复现配置', sz=18, bold=True, clr=DARK)
txt(slide15, 550000, 5150000, 7900000, 1200000, [
    'Base Model：LLaVA-v1.5-7b（Vicuna-7B + CLIP-ViT-Large）',
    'Benchmark：UCIT — 6任务连续学习序列',
    '（ImageNet-R → ArxivQA → VizWiz-Cap → IconQA → CLEVR → Flickr30k）',
    '训练框架：DeepSpeed ZeRO-2，LoRA rank=64，6个专家',
], sz=15, clr=GRAY, sp=Pt(4))

# ============================================================
# 原Slide 17 → 现在是 index 17: 基线复现结果
# ============================================================
slide17 = prs.slides[17]
clear_slide(slide17)
title_bar(slide17, 'HiDe-LLaVA 基线复现结果')

txt(slide17, 500000, 830000, 8100000, 400000,
    '6任务连续学习序列的准确率矩阵（行=训练阶段，列=数据集）',
    sz=16, clr=GRAY)

img(slide17, os.path.join(charm_dir, 'heatmap_hide.png'),
    650000, 1300000, w=7800000, h=5100000)

# ============================================================
# 原Slide 18 → 现在是 index 18: 改进方法
# ============================================================
slide18 = prs.slides[18]
clear_slide(slide18)
title_bar(slide18, '改进方法：引入 Task Prompt')

txt(slide18, 500000, 850000, 8100000, 400000,
    '改进动机', sz=20, bold=True, clr=BLUE)

rbox(slide18, 400000, 1350000, 8300000, 1200000, L_BLUE, BLUE)
txt(slide18, 550000, 1450000, 7900000, 1050000, [
    'HiDe-LLaVA的任务路由仅依赖CLIP特征与锚点的相似度匹配，',
    '缺少对任务身份的显式编码。我们尝试为每个任务引入可学习的',
    'Task Prompt，在输入层面为模型提供额外的任务身份信息。',
], sz=16, clr=DARK, sp=Pt(5))

txt(slide18, 500000, 2800000, 8100000, 400000,
    '具体做法', sz=20, bold=True, clr=BLUE)

items = [
    ('定义', '为每个任务定义可学习提示向量 [10×4096]'),
    ('注入', '在输入序列 system token 之后插入当前任务的 Prompt'),
    ('训练', '采用独立AdamW优化器更新Prompt，与LoRA参数分开训练'),
]
for i, (label, desc) in enumerate(items):
    y = 3300000 + i * 600000
    rbox(slide18, 500000, y, 1400000, 460000, BLUE, BLUE)
    txt(slide18, 580000, y + 90000, 1240000, 350000,
        label, sz=18, bold=True, clr=WHITE)
    txt(slide18, 2100000, y + 90000, 6200000, 350000,
        desc, sz=16, clr=DARK)

rbox(slide18, 400000, 5200000, 8300000, 750000, L_GRAY, GRAY)
txt(slide18, 550000, 5300000, 7900000, 550000, [
    '输入序列:  [System Token] → [Task Prompt (10)] → [Image (576)] → [Text]',
    '标签掩码:  Prompt位置标记为IGNORE_INDEX，不参与损失计算',
], sz=15, clr=GRAY, sp=Pt(5))

txt(slide18, 500000, 6100000, 8100000, 400000,
    '→ 在HiDe-LLaVA基础上，仅增加 10×4096 = 40K 参数/任务',
    sz=16, bold=True, clr=GREEN)

# ============================================================
# 原Slide 19 → 现在是 index 19: 实验结果
# ============================================================
slide19 = prs.slides[19]
clear_slide(slide19)
title_bar(slide19, '实验结果与分析')

txt(slide19, 500000, 830000, 8100000, 350000,
    '各数据集遗忘曲线对比（红: HiDe-LLaVA  蓝: HiDe-LLaVA + Task Prompt）',
    sz=15, clr=GRAY)

img(slide19, os.path.join(charm_dir, 'forgetting_overlay.png'),
    100000, 1150000, w=8900000, h=4250000)

cy = 5500000
cw, ch, gap = 2700000, 650000, 150000

rbox(slide19, 300000, cy, cw, ch, L_GREEN, GREEN)
txt(slide19, 420000, cy + 80000, cw - 240000, ch - 160000, [
    '✓ ArxivQA +1.83%',
    '✓ IconQA +5.03%',
], sz=14, clr=GREEN, sp=Pt(3))

rbox(slide19, 300000 + cw + gap, cy, cw, ch, L_RED, RED)
txt(slide19, 420000 + cw + gap, cy + 80000, cw - 240000, ch - 160000, [
    '✗ CLEVR -7.50%',
    '  Prompt位置干扰位置编码',
], sz=14, clr=RED, sp=Pt(3))

rbox(slide19, 300000 + (cw + gap) * 2, cy, cw, ch, L_BLUE, BLUE)
txt(slide19, 420000 + (cw + gap) * 2, cy + 80000, cw - 240000, ch - 160000, [
    '→ 整体 67.61% vs 67.69%',
    '  需进一步优化架构',
], sz=14, clr=BLUE, sp=Pt(3))

# ============================================================
# 原Slide 20 → 现在是 index 20: 分隔页
# ============================================================
slide20 = prs.slides[20]
for shape in slide20.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.text.strip() == '04':
                for run in p.runs:
                    run.text = '05'
                if not p.runs:
                    p.text = '05'

# ============================================================
# 原Slide 21 → 现在是 index 21: 后续计划
# ============================================================
slide21 = prs.slides[21]
clear_slide(slide21)
title_bar(slide21, '后续工作计划')

txt(slide21, 500000, 850000, 8100000, 350000,
    '基于实验分析，明确了后续改进方向：',
    sz=17, bold=True, clr=DARK)

txt(slide21, 500000, 1200000, 8100000, 350000,
    '待解决问题与对策', sz=19, bold=True, clr=BLUE)

issues = [
    ('Task Prompt架构进一步优化',
     '调整Prompt位置，并改为驱动底层LoRA自适应融合系数，与顶层MoE形成互补',
     L_RED, RED),
    ('缺乏旧任务记忆回顾',
     '引入经验回放，保留5-10%旧任务样本进行混合回放',
     L_ORANGE, ORANGE),
    ('任务序列设定不够真实',
     '构建非固定任务序列，允许任务以不同子集交错重现',
     L_BLUE, BLUE),
]
for i, (problem, solution, bg, border) in enumerate(issues):
    y = 1600000 + i * 700000
    rbox(slide21, 400000, y, 8300000, 580000, bg, border)
    txt(slide21, 550000, y + 60000, 7800000, 240000,
        problem, sz=16, clr=border, bold=True)
    txt(slide21, 550000, y + 320000, 7800000, 240000,
        '→ ' + solution, sz=15, clr=GREEN)

# 非固定任务序列示意
txt(slide21, 500000, 4100000, 8100000, 300000,
    '非固定任务序列示意', sz=16, bold=True, clr=DARK)

# 示意图：用小色块表示任务子集序列
task_colors = [
    (BLUE, 'A-a'), (ORANGE, 'B-b'), (GREEN, 'C-a'),
    (BLUE, 'A-b'), (GREEN, 'C-c'), (PURPLE, 'D-a'),
    (ORANGE, 'B-c'), (RGBColor(0x1A, 0xBC, 0x9C), '...'),
]
bw, bh, bgap = 950000, 450000, 80000
bx_start = 350000
for i, (color, label) in enumerate(task_colors):
    x = bx_start + i * (bw + bgap)
    rbox(slide21, x, 4400000, bw, bh, color, color)
    txt(slide21, x, 4470000, bw, 300000,
        label, sz=14, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)

txt(slide21, 500000, 4950000, 8100000, 300000,
    '同一任务以不同数据子集在序列中多次出现，更贴近真实应用场景',
    sz=14, clr=GRAY)

txt(slide21, 500000, 5350000, 8100000, 350000,
    '进度安排', sz=19, bold=True, clr=BLUE)

phases = [
    ('3.20 — 4.06', '架构改进与验证', L_BLUE, BLUE),
    ('4.07 — 4.20', '对比实验与消融', L_ORANGE, ORANGE),
    ('4.21 — 4.27', '论文初稿', L_PURPLE, PURPLE),
    ('4.28 — 5.15', '完善与答辩', L_GREEN, GREEN),
]
pw, pgap, px_start = 1950000, 100000, 350000
for i, (time, label, bg, border) in enumerate(phases):
    x = px_start + i * (pw + pgap)
    rbox(slide21, x, 5750000, pw, 750000, bg, border)
    txt(slide21, x + 100000, 5800000, pw - 200000, 300000,
        time, sz=13, clr=GRAY)
    txt(slide21, x + 100000, 6080000, pw - 200000, 350000,
        label, sz=15, clr=border, bold=True)

# ============================================================
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Total slides: {len(prs.slides)}')
